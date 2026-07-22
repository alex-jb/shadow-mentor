#!/usr/bin/env python3
"""Shadow capstone deck V2.1 — targeted polish over V2 (same identity, no redesign).

Fixes: fact refresh (1,892/1,895 from feat/shadow-shared-story-adapters @ 5f655e8),
slide-5 stray-text defect, slide-1 tighter arc crop, slide-4 semantic shapes,
slide-6 profile differentiation, slide-7 fixture label, slide-8 tighter crop,
slide-9 honest Unity (text, not a mislabelled Three.js image), slide-10/11 status
refresh, slide-12 projector contrast, backup cleanup. All text from
capstone-facts-v2-1.json (single source). Does not overwrite V2.
"""
import json, os
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
SUB = os.path.dirname(HERE)
FIG = os.path.join(SUB, "figures-v2")       # V2 assets (reused)
FIG1 = os.path.join(SUB, "figures-v2-1")    # V2.1 tighter crops
F = json.load(open(os.path.join(SUB, "capstone-facts-v2-1.json")))
T = F["tests"]
TF = F["canonical_tamper_fixture"]

def C(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
BG      = C("090D12"); PANEL = C("111820"); PANEL2 = C("18212B")
INK     = C("F2F5F7"); INK2  = C("9DA9B5"); MUTE  = C("8A94A0")   # MUTE brightened for projectors
VERIFY  = C("2FD19A"); WARN  = C("F2C14E"); TAMPER = C("FF5F6D"); INFO = C("5CA8FF")
LINE    = C("2A3340"); LINE2 = C("3A4653")
FONT = "Arial"

prs = Presentation(); prs.slide_width = I(13.333); prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = 13.333, 7.5

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), I(SW), I(SH))
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,after=4,line=1.05):
    tb=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[[(runs,18,INK,False)]]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(after); p.line_spacing=line
        if isinstance(para,tuple): para=[para]
        for seg in para:
            t,sz,col,bold=(seg+(False,))[:4] if len(seg)<4 else seg
            r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=FONT
    return tb

def box(s,x,y,w,h,fill=None,ln=None,lnw=1.0,round_=True,adj=0.06):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp=s.shapes.add_shape(shp,I(x),I(y),I(w),I(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb=ln; sp.line.width=Pt(lnw)
    if round_:
        try: sp.adjustments[0]=adj
        except Exception: pass
    return sp

def glyph(s,shape,x,y,w,h,fill,ln=None,lnw=1.5):
    """Draw a semantic MSO_SHAPE, falling back to a rounded rect if unavailable."""
    try: sp=s.shapes.add_shape(shape,I(x),I(y),I(w),I(h))
    except Exception: sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,I(x),I(y),I(w),I(h))
    sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb=ln; sp.line.width=Pt(lnw)
    return sp

def rule(s,x,y,w,col=INFO,h=0.045): box(s,x,y,w,h,fill=col,round_=False)

def title(s,kicker,ttl,kcol=INFO):
    txt(s,0.7,0.5,11.9,0.4,[[(kicker.upper(),13,kcol,True)]])
    txt(s,0.7,0.84,11.9,0.9,[[(ttl,32,INK,True)]])
    rule(s,0.72,1.68,1.4,col=kcol)

def footer(s,n):
    txt(s,0.7,7.05,7,0.3,[[("Alex Ji · Shadow · Capstone I",9,MUTE,False)]])
    txt(s,11.7,7.05,1.0,0.3,[[(f"{n}",9,MUTE,False)]],align=PP_ALIGN.RIGHT)

def pic(s,path,x,y,w=None,h=None,frame=True):
    if not os.path.exists(path):
        box(s,x,y,w or 4,h or 3,fill=PANEL,ln=LINE); return None
    kw={}
    if w: kw["width"]=I(w)
    if h: kw["height"]=I(h)
    im=s.shapes.add_picture(path,I(x),I(y),**kw)
    if frame:
        box(s,Emu(im.left).inches,Emu(im.top).inches,Emu(im.width).inches,Emu(im.height).inches,fill=None,ln=LINE2,lnw=1.0,round_=False)
    return im

def figp(fig,name): return os.path.join(fig,name)
def chip(s,x,y,text,col,w=2.0,sz=11):
    box(s,x,y,w,0.4,fill=None,ln=col,lnw=1.25,adj=0.5)
    txt(s,x,y,w,0.4,[[(text,sz,col,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def callout(s,x,y,num,head,sub,col):
    box(s,x,y,0.42,0.42,fill=col,ln=None,adj=0.5)
    txt(s,x,y,0.42,0.42,[[(str(num),15,BG,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.55,y-0.04,4.6,0.7,[[(head,13,col,True)],[(sub,11.5,INK2,False)]],line=1.1)

# ===== 1 HERO (tighter, larger arc crop) =====
s=slide()
rule(s,0,0,SW,col=INFO,h=0.1)
txt(s,0.9,1.15,11.6,0.4,[[("CAPSTONE IN COMPUTER SCIENCE I",13,INFO,True)]])
txt(s,0.9,1.6,7.2,1.4,[[("Shadow",58,INK,True)]])
txt(s,0.9,2.78,8.2,1.0,[[("A Cryptographically Verifiable Evidence and Spatial",19,INK,False)],
    [("Audit System for AI-Assisted Decisions",19,INK,False)]],line=1.15)
rule(s,0.92,3.95,1.8,col=INFO)
txt(s,0.9,4.12,8.5,0.5,[[("From AI answers to independently verifiable decision evidence.",14,INFO,False)]])
# large cropped arc band (aspect ~4.71 → fits above footer)
im=pic(s,figp(FIG1,"s1-arc-crop.png"),0.9,4.7,w=9.6)
txt(s,0.9,6.82,9.6,0.3,[[("Real render — Shadow 3D audit arc · research prototype",10,MUTE,True)]])
# author block right
txt(s,10.75,4.75,2.4,1.9,[[("Alex Ji",14,INK,True)],[("M.S. Computer Science",11.5,INK2,False)],
    [("Yeshiva University",11.5,INK2,False)],[("",6,INK,False)],[("Practice presentation",11,MUTE,False)],[("2026-07-21",11,MUTE,False)]],line=1.3)

# ===== 2 TRUST GAP =====
s=slide(); title(s,"The trust gap","An answer is not an evidence record")
box(s,0.9,2.3,3.7,1.6,fill=PANEL,ln=LINE,adj=0.12)
txt(s,1.1,2.55,3.3,1.1,[[("AI ANSWER",12,INK2,True)],[("“Approve the loan.”",17,INK,True)],[("fluent · unverifiable",11,MUTE,False)]],line=1.25)
txt(s,4.75,2.3,1.0,1.6,[[("≠",50,TAMPER,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
box(s,5.9,2.3,6.5,1.6,fill=PANEL,ln=VERIFY,lnw=1.25,adj=0.12)
txt(s,6.15,2.5,6.1,1.2,[[("EVIDENCE RECORD",12,VERIFY,True)],
    [("source → tool action → sequence → signature → verifier",14,INK,False)],
    [("structured · signed · independently checkable",11,MUTE,False)]],line=1.3)
for i,(a,b) in enumerate([("SOURCE","What did it rely on?"),("ACTION","What did it do?"),
      ("INTEGRITY","Was it changed after?"),("INDEPENDENCE","Can someone else verify it?")]):
    x=0.9+i*3.0
    rule(s,x,4.55,0.45,col=INFO,h=0.05)
    txt(s,x,4.7,2.75,1.0,[[(a,14,INFO,True)],[(b,13,INK,False)]],line=1.2)
txt(s,0.9,6.2,11.5,0.5,[[("Ordinary logs are mutable and platform-specific; a model’s explanation is generated prose, not a record.",13,INK2,False)]])
footer(s,2)

# ===== 3 EVOLUTION =====
s=slide(); title(s,"Project evolution","From more opinions to verifiable evidence")
for i,(t,d,c) in enumerate([("ORALLEXA","Multiple AI perspectives debate an answer",INFO),
       ("THE LESSON","More voices do not create verifiability",WARN),
       ("SHADOW","Portable evidence + independent verification",VERIFY)]):
    y=2.2+i*1.42
    box(s,0.9,y,0.12,1.05,fill=c,round_=False)
    txt(s,1.25,y,3.4,1.1,[[(t,19,c,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.9,y,7.5,1.1,[[(d,16,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<2: txt(s,1.9,y+0.95,3,0.5,[[("↓",22,MUTE,True)]])
txt(s,0.9,6.5,11.5,0.5,[[("Trust moved ",15,INK2,False),("down a layer",15,VERIFY,True),
    (" — from the voices to the evidence beneath them. A refinement, not a failure.",15,INK2,False)]])
footer(s,3)

# ===== 4 LIFECYCLE — distinct semantic shapes =====
s=slide(); title(s,"How Shadow records evidence","A cryptographically bound sequence")
stages=[("Source","doc",MSO_SHAPE.FLOWCHART_DOCUMENT,INK2,"document"),
        ("Action","tool",MSO_SHAPE.GEAR_6,INK2,"gear"),
        ("Event","seq 2 · a3f9…",MSO_SHAPE.HEXAGON,INFO,"hexagon node"),
        ("Hash chain","commits prev",MSO_SHAPE.DONUT,INFO,"linked rings"),
        ("Ed25519 seal","batch root",MSO_SHAPE.STAR_8_POINT,INFO,"seal"),
        ("Verify","independent",MSO_SHAPE.OVAL,VERIFY,"check")]
n=len(stages); w=1.82; gap=(11.9-n*w)/(n-1)
for i,(lab,sub,shp,c,_) in enumerate(stages):
    x=0.72+i*(w+gap); cx=x+w/2
    # icon
    gs=0.62
    if lab=="Hash chain":
        glyph(s,MSO_SHAPE.DONUT,cx-0.55,2.75,0.5,0.5,None,ln=c,lnw=2.0)
        glyph(s,MSO_SHAPE.DONUT,cx+0.05,2.75,0.5,0.5,None,ln=c,lnw=2.0)
        box(s,cx-0.12,2.97,0.24,0.06,fill=c,round_=False)
    elif lab=="Verify":
        glyph(s,MSO_SHAPE.OVAL,cx-gs/2,2.7,gs,gs,None,ln=c,lnw=2.0)
        txt(s,cx-gs/2,2.68,gs,gs,[[("✓",22,c,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    else:
        fillc = c if shp in (MSO_SHAPE.HEXAGON,MSO_SHAPE.STAR_8_POINT) else None
        g=glyph(s,shp,cx-gs/2,2.7,gs,gs,fillc if fillc else None,ln=(None if fillc else c),lnw=2.0)
    txt(s,x,3.55,w,0.7,[[(lab,13.5,(c if c!=INK2 else INK),True)],[(sub,10.5,INK2,False)]],align=PP_ALIGN.CENTER,line=1.05)
    if i<n-1:
        box(s,x+w+0.02,3.0,gap-0.04,0.05,fill=(INFO if i>=1 else LINE2),round_=False)
txt(s,0.9,4.75,11.5,0.9,[[("Change one earlier event → the ",15,INK,False),("first broken link",15,TAMPER,True),
    (" and its downstream impact become visible.",15,INK,False)]])
chip(s,0.9,5.55,"VERIFIES OFFLINE",VERIFY,w=2.4)
chip(s,3.5,5.55,"SHA-256 + Ed25519",INFO,w=2.6)
footer(s,4)

# ===== 5 ARCHITECTURE — bounded, no stray text, no colliding rail =====
s=slide(); title(s,"System architecture","One core, three profiles, many surfaces")
def band(y,label,items,tag,tagcol):
    box(s,0.9,y,11.55,1.2,fill=PANEL,ln=LINE,lnw=1.0)
    txt(s,1.15,y+0.12,2.5,1.0,[[(label,14.5,INFO,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.5,y+0.14,6.3,1.0,[[(items,12,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,line=1.2)
    # maturity chip INSIDE the band (bounded well within slide)
    box(s,9.95,y+0.38,2.35,0.44,fill=None,ln=tagcol,lnw=1.25,adj=0.5)
    txt(s,9.95,y+0.38,2.35,0.44,[[(tag,10.5,tagcol,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
band(2.05,"EXPERIENCES","CLI · MCP (11) · HTTP · verify.html · Unity Shadow Lens · Three.js replay · 3 explainers","MIXED MATURITY",WARN)
band(3.45,"PROFILES","banking-v1     ·     data-science-v1     ·     coding-agent-v1","ONE GRAMMAR",INFO)
band(4.85,"SHADOW CORE","canonical events · hash chain · Ed25519 · source maps · claim–evidence graph · ingest audit","HOST-TESTED",VERIFY)
txt(s,0.9,6.35,11.55,0.6,[[("A portable ",13,INK2,False),("evidence bundle",13,INFO,True),
    (" flows through every layer — verifiable by anyone, offline. Interface maturity is labelled, not hidden.",13,INK2,False)]])
footer(s,5)

# ===== 6 PROFILES — three visually distinct columns =====
s=slide(); title(s,"Three profiles, one grammar","Same machinery, domain-specific evidence")
colx=[0.9,5.15,9.4]; cw=3.6
# Banking (real crop inset)
x=colx[0]; box(s,x,2.15,cw,3.75,fill=PANEL,ln=LINE,lnw=1.0)
glyph(s,MSO_SHAPE.FLOWCHART_DOCUMENT,x+0.2,2.35,0.55,0.6,None,ln=VERIFY,lnw=2.0)
txt(s,x+0.9,2.4,2.6,0.5,[[("BANKING",15,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,x+0.2,3.05,cw-0.4,0.5,[[("document → risk claim → source quote",12,INK2,False)]],line=1.15)
pic(s,figp(FIG1,"s6-banking-trust.png"),x+0.2,3.6,w=cw-0.4)
txt(s,x+0.2,5.55,cw-0.4,0.3,[[("live audit-chain fixture (real render)",9.5,MUTE,True)]])
# Data Science (concept motif — database → model → metric)
x=colx[1]; box(s,x,2.15,cw,3.75,fill=PANEL,ln=LINE,lnw=1.0)
glyph(s,MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,x+0.2,2.35,0.55,0.6,None,ln=INFO,lnw=2.0)
txt(s,x+0.9,2.4,2.6,0.5,[[("DATA SCIENCE",15,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,x+0.2,3.05,cw-0.4,0.5,[[("dataset → model → metric → selection",12,INK2,False)]],line=1.15)
for j,(a,b) in enumerate([("DATASET","versioned input"),("MODEL","tool/model action"),("METRIC","scored"),("SELECTION","✓ chosen")]):
    yy=3.65+j*0.52
    box(s,x+0.2,yy,cw-0.4,0.44,fill=PANEL2,ln=LINE,lnw=0.75)
    txt(s,x+0.35,yy,2.0,0.44,[[(a,11,INFO,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+1.9,yy,cw-2.1,0.44,[[(b,10.5,INK2,False)]],anchor=MSO_ANCHOR.MIDDLE)
# Coding Agent (concept motif — diff → tests → commit)
x=colx[2]; box(s,x,2.15,cw,3.75,fill=PANEL,ln=LINE,lnw=1.0)
glyph(s,MSO_SHAPE.FLOWCHART_MERGE,x+0.2,2.35,0.55,0.55,None,ln=WARN,lnw=2.0)
txt(s,x+0.9,2.4,2.6,0.5,[[("CODING AGENT",15,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,x+0.2,3.05,cw-0.4,0.5,[[("issue → tool call → diff → tests → commit",12,INK2,False)]],line=1.15)
# mini diff block
box(s,x+0.2,3.65,cw-0.4,1.15,fill=PANEL2,ln=LINE,lnw=0.75)
txt(s,x+0.35,3.72,cw-0.7,1.05,[[("+ ",11,VERIFY,True),("added verified path",10.5,INK2,False)],
    [("+ ",11,VERIFY,True),("guard on ingest",10.5,INK2,False)],
    [("−  ",11,TAMPER,True),("removed silent skip",10.5,INK2,False)]],line=1.35)
box(s,x+0.2,4.95,1.65,0.44,fill=PANEL2,ln=VERIFY,lnw=0.9,adj=0.5)
txt(s,x+0.2,4.95,1.65,0.44,[[("tests ✓✓",11,VERIFY,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
box(s,x+1.95,4.95,cw-2.15,0.44,fill=PANEL2,ln=INFO,lnw=0.9,adj=0.5)
txt(s,x+1.95,4.95,cw-2.15,0.44,[[("commit a3f9…",11,INFO,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# shared grammar strip
box(s,0.9,6.05,12.1,0.72,fill=PANEL2,ln=INFO,lnw=1.0)
txt(s,0.9,6.05,12.1,0.72,[[("Underneath, identical:   sequence → hash chain → signature → source resolution",15,INFO,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
footer(s,6)

# ===== 7 TAMPER (hero) — canonical seq-3 fixture, labelled =====
s=slide()
txt(s,0.7,0.4,11.9,0.4,[[("THE KEY MOMENT",13,TAMPER,True)]])
txt(s,0.7,0.74,7.5,0.7,[[("Exact tamper localization",32,INK,True)]])
rule(s,0.72,1.54,1.4,col=TAMPER)
# fixture label parked top-right, clear of the title
box(s,8.5,0.62,4.1,0.5,fill=None,ln=WARN,lnw=1.0,adj=0.4)
txt(s,8.5,0.62,4.1,0.5,[[("Spatial replay fixture — first failed sequence 3",11.5,WARN,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
pic(s,figp(FIG1,"s7-tamper.png"),0.7,1.9,w=8.0)
callout(s,9.0,2.1,1,"Seq 3 — TAMPERED","the mutated Council-claims node",TAMPER)
callout(s,9.0,3.0,2,"Downstream 4, 5, 6","NOT VERIFIED — chain invalidated",WARN)
callout(s,9.0,3.9,3,"Analytical correctness","NOT EVALUATED — not judged",INFO)
box(s,9.0,4.9,3.6,1.55,fill=PANEL,ln=LINE)
txt(s,9.2,5.03,3.3,1.35,[[("Shadow reports:",12,INK2,True)],
    [("first failed seq: 3",13,TAMPER,True)],
    [("reason: prev_hash_mismatch",12,INK,False)],
    [("downstream: 4, 5, 6",12,WARN,False)]],line=1.3)
txt(s,0.7,6.5,8.0,0.4,[[("Integrity failure does not determine analytical correctness — the two stay independent.",12.5,INK2,True)]])
footer(s,7)

# ===== 8 VERIFY THE VERIFIER — tighter crop =====
s=slide(); title(s,"Verify the Verifier","Can you trust the verifier page itself?")
im=pic(s,figp(FIG1,"s8-verifier-crop.png"),5.35,2.0,w=7.3)
# caption placed clearly BELOW the image frame (no overlap)
imb = 2.0 + (Emu(im.height).inches if im else 3.5)
txt(s,5.35,imb+0.08,7.3,0.3,[[("Real render — Verify-the-Verifier (fixture release key)",10,MUTE,True)]])
callout(s,0.9,2.2,1,"Manifest signature","VERIFIED",VERIFY)
callout(s,0.9,3.1,2,"Assets match","SIGNED MANIFEST",VERIFY)
callout(s,0.9,4.0,3,"Independent comparison","NOT PERFORMED",WARN)
txt(s,0.9,5.05,4.3,1.3,[[("A page hashing itself is ",13,INK,False),("not independent trust",13,TAMPER,True),(".",13,INK,False)],
    [("",6,INK,False)],
    [("Signing today is a ",12.5,INK2,False),("FIXTURE release key",12.5,WARN,True),(", not production.",12.5,INK2,False)]],line=1.3)
footer(s,8)

# ===== 9 SPATIAL — honest Unity (text, not a mislabelled image) =====
s=slide(); title(s,"Spatial audit replay","3D for sequence, provenance, tamper propagation")
pic(s,figp(FIG1,"s9-threejs-arc.png"),0.9,2.0,w=5.75)
txt(s,0.9,4.55,5.75,0.3,[[("THREE.JS — arc layout · browser-rendered",10,MUTE,True)]])
pic(s,figp(FIG1,"s9-threejs-tamper.png"),6.85,2.0,w=5.6)
txt(s,6.85,4.55,5.6,0.3,[[("THREE.JS — tamper cascade · browser-rendered",10,MUTE,True)]])
# status panel (Unity described honestly, not pictured)
box(s,0.9,4.95,7.35,1.65,fill=PANEL,ln=LINE,lnw=1.0)
txt(s,1.1,5.05,7.05,1.5,
    [[("THREE.JS STORY ADAPTER",11.5,INFO,True),("   host-tested contract + browser-rendered · 5 layouts + 2D fallback",10.5,INK2,False)],
     [("SHARED SCENE / STORY CONTRACT",11.5,VERIFY,True),("   authored + host-tested (shadow-3d-scene-v1)",10.5,INK2,False)],
     [("UNITY SHADOW LENS",11.5,WARN,True),("   unity-authored · C# contract-drift host-tested · Android-built",10.5,INK2,False)],
     [("   ",10.5,INK2,False),("device-validation pending · not pictured (no Game View capture)",10.5,MUTE,True)]],line=1.45)
# guardrails reduced to one muted line
txt(s,8.45,4.98,4.5,1.6,[[("Honest bounds:",10.5,INK2,True)],
    [("· not eye tracking",10.5,MUTE,False)],[("· no 6DoF claim on mock",10.5,MUTE,False)],
    [("· no user study yet (RQ4)",10.5,MUTE,False)]],line=1.35)
footer(s,9)

# ===== 10 EVALUATION — refreshed numbers, no stale note =====
s=slide(); title(s,"Evaluation and evidence","Measured today — not inferred")
box(s,0.9,2.15,3.6,3.9,fill=PANEL,ln=VERIFY,lnw=1.25)
txt(s,0.9,2.75,3.6,1.4,[[(f"{T['passed']:,}",46,VERIFY,True)],[(f"of {T['total']:,} host tests",15,INK,False)]],align=PP_ALIGN.CENTER,line=1.0)
txt(s,0.9,4.5,3.6,1.4,[[(f"{T['failed']} failed",15,INK,True)],[(f"{T['skipped']} skipped (env-gated)",13,INK2,False)],
    [(f"{T['suites']} suites",12,MUTE,False)]],align=PP_ALIGN.CENTER,line=1.35)
cats=[("BROWSER","Chromium 149 · EN + 简体中文 · 0 external requests · 0 CSP violations · offline verify",VERIFY),
      ("ANDROID","APK built · 24.4 MB · SHA-256 recorded · NOT device-validated",WARN),
      ("SPATIAL","Scene + story contract host-tested · Three.js adapter host-tested + browser-rendered · Unity C# contract-drift tested · no user study",INFO)]
for i,(a,b,c) in enumerate(cats):
    y=2.15+i*1.32
    box(s,4.8,y,7.65,1.15,fill=PANEL,ln=LINE,lnw=1.0)
    box(s,4.8,y,0.1,1.15,fill=c,round_=False)
    txt(s,5.05,y+0.1,7.25,1.0,[[(a,13,c,True)],[(b,11.5,INK,False)]],line=1.2,anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.9,6.5,11.5,0.3,[[(f"Re-run 2026-07-21 on {T['run_on']}.",10,MUTE,True)]])
footer(s,10)

# ===== 11 NOW vs CAPSTONE II — refreshed =====
s=slide(); title(s,"Current state and Capstone II","Accurate statuses, not blanket ‘pending’")
box(s,0.9,2.15,5.6,4.15,fill=PANEL,ln=VERIFY,lnw=1.0)
txt(s,1.15,2.3,5.1,0.4,[[("NOW",15,VERIFY,True)]])
for i,t in enumerate(["Core evidence + hash-chain verification","Claim–evidence graph (host-tested)",
      "Bilingual offline verifier · 3 explainers","Shared scene + story contract — authored + host-tested",
      "Three.js story adapter — host-tested + browser-rendered","Unity C# contract-drift tested · Android build"]):
    txt(s,1.15,2.85+i*0.55,5.2,0.5,[[("✓  ",13,VERIFY,True),(t,12,INK,False)]])
box(s,6.85,2.15,5.6,4.15,fill=PANEL,ln=WARN,lnw=1.0)
txt(s,7.1,2.3,5.1,0.4,[[("NEXT — CAPSTONE II",15,WARN,True)]])
for i,t in enumerate(["Unity production integration + editmode/playmode","Beam Pro device validation",
      "XREAL native input / camera path","Device performance measurement",
      "Spatial-vs-2D user study (RQ4)","Production signing · KMS/HSM · governance"]):
    txt(s,7.1,2.85+i*0.55,5.2,0.5,[[("○  ",13,WARN,True),(t,12,INK,False)]])
txt(s,0.9,6.5,11.6,0.4,[[("Ingest audit: ",11.5,INK2,True),("structural host-tested",11.5,VERIFY,True),
    (" · ",11.5,INK2,False),("semantic production-evaluation pending",11.5,WARN,True),
    (" — not device-pending.",11.5,INK2,True)]])
footer(s,11)

# ===== 12 CLOSING — projector contrast =====
s=slide()
rule(s,0,0,SW,col=INFO,h=0.1)
txt(s,0.9,1.7,11.5,0.4,[[("CONTRIBUTION",13,INFO,True)]])
txt(s,0.9,2.4,11.6,2.4,[[("Shadow does not ask an auditor to",28,INK,False)],
    [("trust the AI’s answer.",28,INK,False)],[("",10,INK,False)],
    [("It gives them evidence they can",30,INK,True)],[("independently verify.",30,INK,True)]],line=1.18)
rule(s,0.92,5.5,1.8,col=INFO)
txt(s,0.9,5.72,11.6,0.9,[[("Sequence.  Sources.  Signatures.  Exact tamper location.",15,INK2,False)],
    [("Integrity — not correctness.",18,VERIFY,True)]],line=1.4)
txt(s,0.9,6.95,11.6,0.35,[[("github.com/alex-jb/shadow-mentor    ·    ",12.5,INFO,True),("Thank you — questions welcome.",12.5,INK,True)]])

# ===== BACKUP 13 status matrix (refreshed, larger labels) =====
s=slide(); title(s,"Backup · Implementation status","Capability ladder — nothing over-claimed",kcol=MUTE)
rows=[("Evidence schema · hash chain · Ed25519","HOST-TESTED",VERIFY),
      ("Offline browser verifier (EN + zh-CN)","BROWSER-RENDERED / RECORDED",VERIFY),
      ("Verify-the-Verifier manifest","FIXTURE-SIGNED",WARN),
      ("Claim–evidence graph · 3 explainers","HOST-TESTED",VERIFY),
      ("Shared scene + guided-story contract","AUTHORED + HOST-TESTED",VERIFY),
      ("Three.js story adapter","HOST-TESTED · BROWSER-RENDERED",INFO),
      ("Unity Shadow Lens","UNITY-AUTHORED · C# DRIFT-TESTED · DEVICE-PENDING",WARN),
      ("Android mock APK","ANDROID-BUILT (not device-validated)",WARN),
      ("Ingest audit — structural / semantic","HOST-TESTED / PRODUCTION-PENDING",INFO),
      ("Spatial comprehension benefit (RQ4)","NOT IMPLEMENTED — no user study",MUTE)]
for i,(a,st,c) in enumerate(rows):
    y=2.0+i*0.5
    txt(s,0.9,y,7.0,0.45,[[(a,13,INK,False)]])
    txt(s,8.1,y,4.7,0.45,[[(st,12.5,c,True)]])
footer(s,13)

# ===== BACKUP 14 trust boundary =====
s=slide(); title(s,"Backup · Cryptographic trust boundary","What a valid signature does and does not prove",kcol=MUTE)
box(s,0.9,2.15,5.6,3.9,fill=PANEL,ln=VERIFY,lnw=1.0)
txt(s,1.15,2.35,5.1,0.4,[[("ESTABLISHES",14,VERIFY,True)]])
for i,t in enumerate(["The bytes match the sealed record","The chain is unbroken (or where it breaks)","The batch root was signed by a known key","Conformance to a domain profile"]):
    txt(s,1.15,2.95+i*0.62,5.1,0.55,[[("✓  ",13,VERIFY,True),(t,13,INK,False)]])
box(s,6.85,2.15,5.6,3.9,fill=PANEL,ln=TAMPER,lnw=1.0)
txt(s,7.1,2.35,5.1,0.4,[[("DOES NOT ESTABLISH",14,TAMPER,True)]])
for i,t in enumerate(["That the cited source was truthful","That the conclusion is correct","That signing is production-grade","That any regulation mandates it"]):
    txt(s,7.1,2.95+i*0.62,5.1,0.55,[[("✗  ",13,TAMPER,True),(t,13,INK,False)]])
footer(s,14)

# ===== BACKUP 15 limitations (refreshed) =====
s=slide(); title(s,"Backup · Limitations","Stated directly",kcol=MUTE)
for i,t in enumerate(["Evidence is fixture data; limited live-provider validation",
      "Production signing not implemented (fixture release key)",
      "Unity is authored + C# contract-drift-tested — not editmode/playmode-run, not device-validated",
      "Device validation pending (Beam Pro); XREAL native not integrated",
      "No completed user study — RQ4 spatial benefit unproven",
      "Semantic ingest audit: production evaluation pending (structural is host-tested)",
      "No proof of source truth; no proof of analytical correctness",
      "No production KMS/HSM, key rotation, or durable storage"]):
    txt(s,0.9,2.05+i*0.55,11.5,0.5,[[("—  ",13,WARN,True),(t,13.5,INK,False)]])
footer(s,15)

# ===== BACKUP 16 demo fallback (safe path) =====
s=slide(); title(s,"Backup · Demo fallback plan","Never depend on a live device",kcol=MUTE)
levels=[("A · Live app","serve verifier package locally",VERIFY),
        ("B · Browser package","open the frozen offline package",INFO),
        ("C · Recorded video",F["demo_fallback_video_label"],WARN),
        ("D · Slide screenshots","embedded figures — one keystroke away",MUTE)]
for i,(a,b,c) in enumerate(levels):
    y=2.2+i*1.0
    box(s,0.9,y,11.5,0.82,fill=PANEL,ln=LINE)
    box(s,0.9,y,0.1,0.82,fill=c,round_=False)
    txt(s,1.2,y,3.0,0.82,[[(a,15,c,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.4,y,7.8,0.82,[[(b,14,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.9,6.5,11.5,0.4,[[("Decision rule: if the live app doesn’t come up cleanly in 20 seconds, open the bookmarked fallback video and keep talking.",12,INK2,True)]])
footer(s,16)

out=os.path.join(HERE,"SHADOW_CAPSTONE_PRACTICE_PRESENTATION_V2_1.pptx")
prs.save(out)
print(f"saved {out} · {len(prs.slides._sldIdLst)} slides")
