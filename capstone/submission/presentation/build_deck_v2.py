#!/usr/bin/env python3
"""Shadow capstone deck V2 — dark forensic identity, real product images.

Design system taken from the real product tokens
(apps/shadow-lens/web/spatial-agent/src/design-tokens.ts):
  bg #090D12 · verified #2FD19A · tampered #FF5F6D · warning #F2C14E ·
  info #5CA8FF · neutral #9DA9B5.
Real screenshots (verifier, Three.js layouts, tamper render) do the visual
work; vector shapes only for architecture / lifecycle / callouts. Every value
comes from capstone-facts-v2.json (single source of truth).

Run:  python3 build_deck_v2.py  → SHADOW_CAPSTONE_PRACTICE_PRESENTATION_V2.pptx
"""
import json, os
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
SUB = os.path.dirname(HERE)
FIG = os.path.join(SUB, "figures-v2")
F = json.load(open(os.path.join(SUB, "capstone-facts-v2.json")))
T = F["tests"]

def C(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
BG      = C("090D12"); PANEL = C("111820"); PANEL2 = C("18212B")
INK     = C("F2F5F7"); INK2  = C("9DA9B5"); MUTE  = C("6B7480")
VERIFY  = C("2FD19A"); WARN  = C("F2C14E"); TAMPER = C("FF5F6D"); INFO = C("5CA8FF")
LINE    = C("2A3340"); LINE2 = C("3A4653")
FONT = "Arial"

prs = Presentation(); prs.slide_width = I(13.333); prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = 13.333, 7.5

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), I(SW), I(SH))
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,after=4,line=1.05):
    tb=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[[(runs,18,INK,False)]]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(after); p.line_spacing=line
        if isinstance(para,tuple): para=[para]
        for seg in para:
            t,sz,col,bold=(seg+(False,))[:4] if len(seg)<4 else seg
            r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=FONT
    return tb

def box(s,x,y,w,h,fill=None,ln=None,lnw=1.0,round_=True,adj=0.06):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp=s.shapes.add_shape(shp,I(x),I(y),I(w),I(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb=ln; sp.line.width=Pt(lnw)
    if round_:
        try: sp.adjustments[0]=adj
        except Exception: pass
    return sp

def rule(s,x,y,w,col=INFO,h=0.045): box(s,x,y,w,h,fill=col,round_=False)

def title(s,kicker,ttl,kcol=INFO):
    txt(s,0.7,0.5,11.9,0.4,[[(kicker.upper(),13,kcol,True)]])
    txt(s,0.7,0.84,11.9,0.9,[[(ttl,32,INK,True)]])
    rule(s,0.72,1.68,1.4,col=kcol)

def footer(s,n,dark=True):
    txt(s,0.7,7.05,7,0.3,[[("Alex Ji · Shadow · Capstone I",9,MUTE,False)]])
    txt(s,11.7,7.05,1.0,0.3,[[(f"{n}",9,MUTE,False)]],align=PP_ALIGN.RIGHT)

def pic(s,name,x,y,w=None,h=None,frame=True):
    p=os.path.join(FIG,name)
    if not os.path.exists(p):
        box(s,x,y,w or 4,h or 3,fill=PANEL,ln=LINE); return None
    kw={}
    if w: kw["width"]=I(w)
    if h: kw["height"]=I(h)
    im=s.shapes.add_picture(p,I(x),I(y),**kw)
    if frame:
        b=box(s,Emu(im.left).inches,Emu(im.top).inches,Emu(im.width).inches,Emu(im.height).inches,fill=None,ln=LINE2,lnw=1.0,round_=False)
    return im

def chip(s,x,y,text,col,w=2.0):
    box(s,x,y,w,0.4,fill=None,ln=col,lnw=1.25,adj=0.5)
    txt(s,x,y,w,0.4,[[(text,11,col,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def callout(s,x,y,num,head,sub,col):
    c=box(s,x,y,0.42,0.42,fill=col,ln=None,adj=0.5)
    txt(s,x,y,0.42,0.42,[[(str(num),15,BG,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.55,y-0.04,4.6,0.7,[[(head,13,col,True)],[(sub,11.5,INK2,False)]],line=1.1)

# ===== 1 HERO =====
s=slide()
rule(s,0,0,SW,col=INFO,h=0.1)
txt(s,0.9,1.4,6.6,0.4,[[("CAPSTONE IN COMPUTER SCIENCE I",13,INFO,True)]])
txt(s,0.9,1.95,6.6,1.4,[[("Shadow",58,INK,True)]])
txt(s,0.9,3.15,6.6,1.4,[[("A Cryptographically Verifiable Evidence",21,INK,False)],
    [("and Spatial Audit System for AI Decisions",21,INK,False)]],line=1.15)
rule(s,0.92,4.5,1.8,col=INFO)
txt(s,0.9,4.72,6.6,0.6,[[("From AI answers to independently verifiable decision evidence.",14,INFO,False)]])
txt(s,0.9,5.75,6.6,1.0,[[("Alex Ji",15,INK,True)],
    [("M.S. Computer Science · Yeshiva University",13,INK2,False)],
    [("Practice presentation · 2026-07-21",12,MUTE,False)]],line=1.35)
pic(s,"v2-3d-current-arc.png",7.35,1.5,w=5.5)
txt(s,7.35,6.35,5.5,0.35,[[("Real render — Shadow 3D audit arc (research prototype)",10,MUTE,True)]])

# ===== 2 TRUST GAP =====
s=slide(); title(s,"The trust gap","An answer is not an evidence record")
# AI bubble
b=box(s,0.9,2.3,3.7,1.6,fill=PANEL,ln=LINE,adj=0.12)
txt(s,1.1,2.55,3.3,1.1,[[("AI ANSWER",12,INK2,True)],[("“Approve the loan.”",17,INK,True)],[("fluent · unverifiable",11,MUTE,False)]],line=1.25)
txt(s,4.75,2.3,1.0,1.6,[[("≠",50,TAMPER,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
b2=box(s,5.9,2.3,6.5,1.6,fill=PANEL,ln=VERIFY,lnw=1.25,adj=0.12)
txt(s,6.15,2.5,6.1,1.2,[[("EVIDENCE RECORD",12,VERIFY,True)],
    [("source → tool action → sequence → signature → verifier",14,INK,False)],
    [("structured · signed · independently checkable",11,MUTE,False)]],line=1.3)
gaps=[("SOURCE","What did it rely on?"),("ACTION","What did it do?"),
      ("INTEGRITY","Was it changed after?"),("INDEPENDENCE","Can someone else verify it?")]
for i,(a,b) in enumerate(gaps):
    x=0.9+i*3.0
    rule(s,x,4.55,0.45,col=INFO,h=0.05)
    txt(s,x,4.7,2.75,1.0,[[(a,14,INFO,True)],[(b,13,INK,False)]],line=1.2)
txt(s,0.9,6.2,11.5,0.5,[[("Ordinary logs are mutable and platform-specific; a model’s explanation is generated prose, not a record.",13,INK2,False)]])
footer(s,2)

# ===== 3 EVOLUTION =====
s=slide(); title(s,"Project evolution","From more opinions to verifiable evidence")
steps=[("ORALLEXA","Multiple AI perspectives debate an answer",INFO),
       ("THE LESSON","More voices do not create verifiability",WARN),
       ("SHADOW","Portable evidence + independent verification",VERIFY)]
for i,(t,d,c) in enumerate(steps):
    y=2.2+i*1.42
    box(s,0.9,y,0.12,1.05,fill=c,round_=False)
    txt(s,1.25,y,3.4,1.1,[[(t,19,c,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.9,y,7.5,1.1,[[(d,16,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if i<2: txt(s,1.9,y+0.95,3,0.5,[[("↓",22,MUTE,True)]])
txt(s,0.9,6.5,11.5,0.5,[[("Trust moved ",15,INK2,False),("down a layer",15,VERIFY,True),
    (" — from the voices to the evidence beneath them. A refinement, not a failure.",15,INK2,False)]])
footer(s,3)

# ===== 4 LIFECYCLE =====
s=slide(); title(s,"How Shadow records evidence","A cryptographically bound sequence")
nodes=[("Source","doc / dataset",INK2),("Action","tool / model",INK2),
       ("Event","seq 2 · a3f9…",INFO),("Hash chain","commits prev",INFO),
       ("Ed25519 seal","batch root",INFO),("Verify","independent",VERIFY)]
n=len(nodes); w=1.82; gap=(11.9-n*w)/(n-1)
for i,(lab,sub,c) in enumerate(nodes):
    x=0.72+i*(w+gap)
    fill = PANEL2 if c==INFO else (PANEL if c==INK2 else PANEL)
    ln = c if c!=INK2 else LINE
    box(s,x,3.0,w,1.4,fill=fill,ln=ln,lnw=1.4)
    txt(s,x,3.12,w,1.2,[[(lab,13.5,(c if c!=INK2 else INK),True)],[(sub,11,INK2,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line=1.1)
    if i<n-1:
        cy=3.7
        cl=box(s,x+w+0.02,cy-0.02,gap-0.04,0.05,fill=(INFO if i>=1 else LINE2),round_=False)
txt(s,0.9,4.95,11.5,1.0,[[("Change one earlier event → the ",15,INK,False),("first broken link",15,TAMPER,True),
    (" and its downstream impact become visible.",15,INK,False)]])
chip(s,0.9,5.7,"VERIFIES OFFLINE",VERIFY,w=2.4)
chip(s,3.5,5.7,"SHA-256 + Ed25519",INFO,w=2.6)
footer(s,4)

# ===== 5 ARCHITECTURE =====
s=slide(); title(s,"System architecture","One core, three profiles, many surfaces")
def band(y,label,items,tags):
    box(s,0.9,y,9.2,1.28,fill=PANEL,ln=LINE,lnw=1.0)
    txt(s,1.15,y+0.14,2.4,1.0,[[(label,15,INFO,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,3.5,y+0.16,6.4,1.0,[[(items,12.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,line=1.2)
    for j,(tg,c) in enumerate(tags):
        txt(s,10.35,y+0.12+j*0.36,2.7,0.35,[[("● ",10,c,True),(tg,10.5,INK2,False)]])
band(2.05,"EXPERIENCES","CLI · MCP (11) · HTTP · verify.html · Unity Shadow Lens · Three.js replay · explainers",
     [("verifier BROWSER-RENDERED",VERIFY),("Unity AUTHORED · device-pending",WARN)])
band(3.5,"PROFILES","banking-v1     ·     data-science-v1     ·     coding-agent-v1",
     [("one verification grammar",INFO)])
band(4.95,"SHADOW CORE","canonical events · hash chain · Ed25519 · source maps · claim–evidence graph · ingest audit (structural)",
     [("HOST-TESTED",VERIFY)])
# side rail
box(s,10.35,2.05,2.05,4.18,fill=PANEL2,ln=INFO,lnw=1.0)
txt(s,10.35,2.05,2.05,4.18,[[("PORTABLE",12,INFO,True)],[("EVIDENCE",12,INFO,True)],[("BUNDLE",12,INFO,True)],[("",8,INK,False)],[("verified by",10,INK2,False)],[("anyone,",10,INK2,False)],[("offline",10,INK2,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line=1.3)
txt(s,0.9,6.5,9.2,0.4,[[("Same verification grammar flows through every layer; interface maturity is labelled, not hidden.",12,INK2,False)]])
footer(s,5)

# ===== 6 PROFILES =====
s=slide(); title(s,"Three profiles, one grammar","Same machinery across domains")
cols=[("Banking","document → risk claim → source quote"),
      ("Data Science","dataset → model → metric → selection"),
      ("Coding Agent","issue → tool call → diff → tests → commit")]
for i,(a,b) in enumerate(cols):
    x=0.9+i*3.95
    box(s,x,2.2,3.6,2.3,fill=PANEL,ln=LINE,lnw=1.0)
    txt(s,x+0.05,2.35,3.5,0.4,[[(f"0{i+1}",22,INFO,True)]],align=PP_ALIGN.CENTER)
    txt(s,x+0.1,2.95,3.4,0.5,[[(a,18,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,x+0.25,3.6,3.1,0.8,[[(b,13,INK2,False)]],align=PP_ALIGN.CENTER,line=1.25)
box(s,0.9,4.95,11.55,1.3,fill=PANEL2,ln=INFO,lnw=1.0)
txt(s,0.9,5.05,11.55,0.5,[[("Underneath, identical:",14,INK2,False)]],align=PP_ALIGN.CENTER)
txt(s,0.9,5.45,11.55,0.6,[[("sequence  →  hash chain  →  signature  →  source resolution",18,INFO,True)]],align=PP_ALIGN.CENTER)
footer(s,6)

# ===== 7 TAMPER (hero, real image) =====
s=slide()
txt(s,0.7,0.45,11.9,0.4,[[("THE KEY MOMENT",13,TAMPER,True)]])
txt(s,0.7,0.8,11.9,0.8,[[("Exact tamper localization",32,INK,True)]])
rule(s,0.72,1.62,1.4,col=TAMPER)
pic(s,"v2-3d-tamper-propagation.png",0.7,1.95,w=8.0)
callout(s,9.0,2.15,1,"Seq 3 — TAMPERED","the mutated Council-claims node",TAMPER)
callout(s,9.0,3.05,2,"Downstream 4, 5, 6","NOT VERIFIED — chain invalidated",WARN)
callout(s,9.0,3.95,3,"Analytical correctness","NOT EVALUATED — not judged",INFO)
box(s,9.0,4.95,3.6,1.55,fill=PANEL,ln=LINE)
txt(s,9.2,5.08,3.3,1.35,[[("Shadow reports:",12,INK2,True)],
    [("first failed seq: 3",13,TAMPER,True)],
    [("reason: prev_hash_mismatch",12,INK,False)],
    [("downstream: 4, 5, 6",12,WARN,False)]],line=1.3)
txt(s,0.7,6.55,8.0,0.4,[[("Integrity failure does not determine analytical correctness — the two stay independent.",12.5,INK2,True)]])
footer(s,7)

# ===== 8 VERIFY THE VERIFIER (real image, big) =====
s=slide(); title(s,"Verify the Verifier","Can you trust the verifier page itself?")
pic(s,"v2-verifier-valid-en.png",5.55,1.95,w=7.1)
txt(s,5.55,6.35,7.1,0.35,[[("Real render — Verify-the-Verifier (fixture release key)",10,MUTE,True)]])
callout(s,0.9,2.2,1,"Manifest signature","VERIFIED",VERIFY)
callout(s,0.9,3.1,2,"Assets match","SIGNED MANIFEST",VERIFY)
callout(s,0.9,4.0,3,"Independent comparison","NOT PERFORMED",WARN)
txt(s,0.9,5.05,4.4,1.2,[[("A page hashing itself is ",13,INK,False),("not independent trust",13,TAMPER,True),(".",13,INK,False)],
    [("",6,INK,False)],
    [("Signing today is a ",12.5,INK2,False),("FIXTURE release key",12.5,WARN,True),(", not production.",12.5,INK2,False)]],line=1.3)
footer(s,8)

# ===== 9 SPATIAL (real Three.js layouts) =====
s=slide(); title(s,"Spatial audit replay","3D for sequence, provenance, tamper propagation")
pic(s,"v2-3d-current-arc.png",0.9,2.0,w=5.75)
txt(s,0.9,4.55,5.75,0.3,[[("Arc layout — audit sequence",10,MUTE,True)]])
pic(s,"v2-3d-layered-dag.png",6.85,2.0,w=5.6)
txt(s,6.85,4.55,5.6,0.3,[[("Layered DAG — provenance",10,MUTE,True)]])
txt(s,0.9,4.95,7.2,1.3,
    [[("THREE.JS REPLAY",12,INFO,True),("   browser-rendered · research prototype · 4 layouts + 2D fallback",11,INK2,False)],
     [("UNITY SHADOW LENS",12,WARN,True),("   unity-authored · Android-built · device validation pending",11,INK2,False)],
     [("SHARED 3D SCENE CONTRACT",12,VERIFY,True),("   authored + host-tested (shadow-3d-scene-v1)",11,INK2,False)]],line=1.5)
for i,t in enumerate(["NOT EYE TRACKING","NO 6DoF ON MOCK","NO USER STUDY YET"]):
    chip(s,8.5+ (i%1)*0, 4.95+i*0.5, t, MUTE, w=3.9)
footer(s,9)

# ===== 10 EVALUATION =====
s=slide(); title(s,"Evaluation and evidence","Measured today — not inferred")
box(s,0.9,2.15,3.6,3.9,fill=PANEL,ln=VERIFY,lnw=1.25)
txt(s,0.9,2.7,3.6,1.5,[[(f"{T['passed']:,}",46,VERIFY,True)],[(f"of {T['total']:,} host tests",15,INK,False)]],align=PP_ALIGN.CENTER,line=1.0)
txt(s,0.9,4.55,3.6,1.3,[[(f"{T['failed']} failed · {T['skipped']} skipped",13,INK2,True)],
    [("re-run 2026-07-21",12,MUTE,False)],[("(V1 said 1,824/1,827)",11,MUTE,False)]],align=PP_ALIGN.CENTER,line=1.3)
cats=[("BROWSER","Chromium 149 · EN + 简体中文 · 0 external requests · 0 CSP violations · offline verify",VERIFY),
      ("ANDROID","APK built · 24.4 MB · SHA-256 recorded · NOT device-validated",WARN),
      ("SPATIAL","Three.js 4 layouts rendered + recorded · scene contract host-tested · no user study",INFO)]
for i,(a,b,c) in enumerate(cats):
    y=2.15+i*1.32
    box(s,4.8,y,7.65,1.15,fill=PANEL,ln=LINE,lnw=1.0)
    box(s,4.8,y,0.1,1.15,fill=c,round_=False)
    txt(s,5.05,y+0.12,7.2,1.0,[[(a,13,c,True)],[(b,12.5,INK,False)]],line=1.25,anchor=MSO_ANCHOR.MIDDLE)
footer(s,10)

# ===== 11 NOW vs CAPSTONE II =====
s=slide(); title(s,"Current state and Capstone II","Accurate statuses, not blanket ‘pending’")
box(s,0.9,2.15,5.6,4.15,fill=PANEL,ln=VERIFY,lnw=1.0)
txt(s,1.15,2.3,5.1,0.4,[[("NOW",15,VERIFY,True)]])
for i,t in enumerate(["Core evidence + hash-chain verification","Claim–evidence graph (host-tested)",
      "Bilingual offline verifier","Three explainers (audit-chain, reason-code, persona)",
      "3D scene contract — authored + host-tested","Three.js 4-layout prototype · Android build"]):
    txt(s,1.15,2.85+i*0.55,5.15,0.5,[[("✓  ",13,VERIFY,True),(t,12.5,INK,False)]])
box(s,6.85,2.15,5.6,4.15,fill=PANEL,ln=WARN,lnw=1.0)
txt(s,7.1,2.3,5.1,0.4,[[("NEXT — CAPSTONE II",15,WARN,True)]])
for i,t in enumerate(["Unity contract integration (pending)","Beam Pro device validation","XREAL native input / camera path",
      "Device performance measurement","Spatial-vs-2D user study (RQ4)","Production signing · KMS/HSM · governance"]):
    txt(s,7.1,2.85+i*0.55,5.15,0.5,[[("○  ",13,WARN,True),(t,12.5,INK,False)]])
txt(s,0.9,6.5,11.5,0.4,[[("Ingest audit: structural host-tested · semantic production evaluation pending — stated separately, not collapsed.",11.5,INK2,True)]])
footer(s,11)

# ===== 12 CLOSING =====
s=slide()
rule(s,0,0,SW,col=INFO,h=0.1)
txt(s,0.9,1.7,11.5,0.4,[[("CONTRIBUTION",13,INFO,True)]])
txt(s,0.9,2.4,11.6,2.4,[[("Shadow does not ask an auditor to",28,INK2,False)],
    [("trust the AI’s answer.",28,INK2,False)],[("",10,INK,False)],
    [("It gives them evidence they can",30,INK,True)],[("independently verify.",30,INK,True)]],line=1.18)
rule(s,0.92,5.5,1.8,col=INFO)
txt(s,0.9,5.75,11.6,0.9,[[("Sequence.  Sources.  Signatures.  Exact tamper location.",15,INK2,False)],
    [("Integrity — not correctness.",18,VERIFY,True)]],line=1.4)
txt(s,0.9,6.95,11.6,0.3,[[("github.com/alex-jb/shadow-mentor    ·    Thank you — questions welcome.",12,MUTE,False)]])

# ===== BACKUP 13 status matrix =====
s=slide(); title(s,"Backup · Implementation status","Capability ladder — nothing over-claimed",kcol=MUTE)
rows=[("Evidence schema · hash chain · Ed25519","HOST-TESTED",VERIFY),
      ("Offline browser verifier (EN + zh-CN)","BROWSER-RENDERED / RECORDED",VERIFY),
      ("Verify-the-Verifier manifest","FIXTURE-SIGNED",WARN),
      ("Claim–evidence graph","HOST-TESTED",VERIFY),
      ("3D scene contract (shadow-3d-scene-v1)","AUTHORED + HOST-TESTED",VERIFY),
      ("Three.js 4-layout replay","BROWSER-RENDERED · RESEARCH-PROTOTYPE",INFO),
      ("Unity Shadow Lens","UNITY-AUTHORED · DEVICE-PENDING",WARN),
      ("Android mock APK","ANDROID-BUILT (not device-validated)",WARN),
      ("Ingest audit — structural / semantic","HOST-TESTED / SOURCE-AUTHORED",INFO),
      ("Spatial comprehension benefit (RQ4)","NOT IMPLEMENTED — no user study",MUTE)]
for i,(a,st,c) in enumerate(rows):
    y=2.0+i*0.5
    txt(s,0.9,y,7.3,0.45,[[(a,13,INK,False)]])
    txt(s,8.3,y,4.3,0.45,[[(st,12,c,True)]])
footer(s,13)

# ===== BACKUP 14 trust boundary =====
s=slide(); title(s,"Backup · Cryptographic trust boundary","What a valid signature does and does not prove",kcol=MUTE)
box(s,0.9,2.15,5.6,3.9,fill=PANEL,ln=VERIFY,lnw=1.0)
txt(s,1.15,2.35,5.1,0.4,[[("ESTABLISHES",14,VERIFY,True)]])
for i,t in enumerate(["The bytes match the sealed record","The chain is unbroken (or where it breaks)","The batch root was signed by a known key","Conformance to a domain profile"]):
    txt(s,1.15,2.95+i*0.62,5.1,0.55,[[("✓  ",13,VERIFY,True),(t,13,INK,False)]])
box(s,6.85,2.15,5.6,3.9,fill=PANEL,ln=TAMPER,lnw=1.0)
txt(s,7.1,2.35,5.1,0.4,[[("DOES NOT ESTABLISH",14,TAMPER,True)]])
for i,t in enumerate(["That the cited source was truthful","That the conclusion is correct","That signing is production-grade","That any regulation mandates it"]):
    txt(s,7.1,2.95+i*0.62,5.1,0.55,[[("✗  ",13,TAMPER,True),(t,13,INK,False)]])
footer(s,14)

# ===== BACKUP 15 limitations =====
s=slide(); title(s,"Backup · Limitations","Stated directly",kcol=MUTE)
lims=["Evidence is fixture data; limited live-provider validation",
      "Production signing not implemented (fixture release key)",
      "Device validation pending (Beam Pro); XREAL native not integrated",
      "No completed user study — RQ4 spatial benefit unproven",
      "No proof of source truth; no proof of analytical correctness",
      "No production KMS/HSM, key rotation, or durable storage",
      "No full PII-retention framework or incident-response process",
      "Semantic ingest audit: production evaluation pending"]
for i,t in enumerate(lims):
    txt(s,0.9,2.05+i*0.55,11.5,0.5,[[("—  ",13,WARN,True),(t,14,INK,False)]])
footer(s,15)

# ===== BACKUP 16 demo fallback =====
s=slide(); title(s,"Backup · Demo fallback plan","Never depend on a live device",kcol=MUTE)
levels=[("A · Live app","serve verifier package locally",VERIFY),
        ("B · Browser package","open the frozen offline package",INFO),
        ("C · Recorded video","media/.../shadow-verify-full-demo.mp4",WARN),
        ("D · Slide screenshots","embedded figures — one keystroke away",MUTE)]
for i,(a,b,c) in enumerate(levels):
    y=2.2+i*1.0
    box(s,0.9,y,11.5,0.82,fill=PANEL,ln=LINE)
    box(s,0.9,y,0.1,0.82,fill=c,round_=False)
    txt(s,1.2,y,3.0,0.82,[[(a,15,c,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,4.4,y,7.8,0.82,[[(b,14,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
txt(s,0.9,6.5,11.5,0.4,[[("Decision rule: if the live app doesn’t come up cleanly in 20 seconds, drop to the recorded video and keep talking.",12,INK2,True)]])
footer(s,16)

out=os.path.join(HERE,"SHADOW_CAPSTONE_PRACTICE_PRESENTATION_V2.pptx")
prs.save(out)
print(f"saved {out} · {len(prs.slides._sldIdLst)} slides")
