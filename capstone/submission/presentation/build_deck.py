#!/usr/bin/env python3
"""Hand-designed Shadow capstone deck (python-pptx).

Not a pandoc dump — every slide is laid out with a consistent design system:
one accent, semantic green/red for the integrity story, generous whitespace,
projector-safe light background, real screenshots where they add credibility.
Content is the same truth as SLIDE_SOURCE.md (the text source of record).

Run:  python3 build_deck.py   →  SHADOW_CAPSTONE_PRACTICE_PRESENTATION.pptx
"""
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "..", "figures")

# ---- design system ----------------------------------------------------------
BG      = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1A, 0x1A, 0x2E)
INK2    = RGBColor(0x55, 0x55, 0x6E)
MUTE    = RGBColor(0x8A, 0x8A, 0x99)
ACC     = RGBColor(0x1A, 0x4D, 0x8F)   # deep blue
ACC_LT  = RGBColor(0xEA, 0xF1, 0xF9)
GREEN   = RGBColor(0x0A, 0x7D, 0x43)
GREEN_LT= RGBColor(0xE7, 0xF3, 0xEC)
RED     = RGBColor(0xC0, 0x39, 0x2B)
RED_LT  = RGBColor(0xFB, 0xEB, 0xE9)
AMBER   = RGBColor(0xB7, 0x79, 0x1F)
AMBER_LT= RGBColor(0xFB, 0xF3, 0xE4)
LINE    = RGBColor(0xDD, 0xDD, 0xD6)
DARK    = RGBColor(0x14, 0x1B, 0x2D)   # title-slide band
FONT    = "Arial"

prs = Presentation()
prs.slide_width  = I(13.333)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), I(SW), I(SH))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def _noshadow(sp): sp.shadow.inherit = False

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line=1.05):
    tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [[(runs, 18, INK, False)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line
        if isinstance(para, tuple): para = [para]
        for seg in para:
            t, sz, col, bold = (seg + (False,))[:4] if len(seg) < 4 else seg
            r = p.add_run(); r.text = t; f = r.font
            f.size = Pt(sz); f.color.rgb = col; f.bold = bold; f.name = FONT
    return tb

def box(s, x, y, w, h, fill=None, ln=None, lnw=1.0, round_=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, I(x), I(y), I(w), I(h)); _noshadow(sp)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb = ln; sp.line.width = Pt(lnw)
    if round_:
        try: sp.adjustments[0] = 0.08
        except Exception: pass
    return sp

def rule(s, x, y, w, col=ACC, h=0.045):
    box(s, x, y, w, h, fill=col, round_=False)

def title_bar(s, kicker, title):
    txt(s, 0.7, 0.5, 11.9, 0.4, [[(kicker.upper(), 12, ACC, True)]])
    txt(s, 0.7, 0.82, 11.9, 0.9, [[(title, 30, INK, True)]])
    rule(s, 0.72, 1.62, 1.5)

def footer(s, n):
    txt(s, 0.7, 7.02, 6, 0.3, [[("Shadow · Capstone I · Alex Ji", 9, MUTE, False)]])
    txt(s, 11.6, 7.02, 1.1, 0.3, [[(f"{n} / 12", 9, MUTE, False)]], align=PP_ALIGN.RIGHT)

def chip(s, x, y, w, text, col, lt):
    box(s, x, y, w, 0.42, fill=lt, ln=col, lnw=1.25)
    txt(s, x, y, w, 0.42, [[(text, 12, col, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================ SLIDE 1 — Title ================================
s = slide()
box(s, 0, 0, SW, SH, fill=DARK, round_=False)
rule(s, 0, 0, SW, col=ACC, h=0.12)
txt(s, 0.9, 2.0, 11.5, 0.5, [[("CAPSTONE IN COMPUTER SCIENCE I", 14, RGBColor(0x9F,0xB6,0xD6), True)]])
txt(s, 0.9, 2.55, 11.6, 1.7, [[("Shadow", 60, RGBColor(0xFF,0xFF,0xFF), True)]])
txt(s, 0.9, 3.75, 11.6, 1.0,
    [[("A Cryptographically Verifiable Evidence and Spatial", 24, RGBColor(0xE7,0xEC,0xF5), False)],
     [("Audit System for AI-Assisted Decisions", 24, RGBColor(0xE7,0xEC,0xF5), False)]], line=1.1)
rule(s, 0.92, 5.05, 2.0, col=ACC)
txt(s, 0.9, 5.3, 11.6, 1.2,
    [[("From Multi-Agent Answers to Independently Verifiable AI Decision Evidence", 15, RGBColor(0x9F,0xB6,0xD6), False)],
     [("Alex Ji   ·   M.S. Computer Science   ·   Yeshiva University", 15, RGBColor(0xE7,0xEC,0xF5), True)],
     [("Practice presentation · 2026-07-21", 13, RGBColor(0x8A,0x9C,0xBE), False)]], line=1.3)

# ============================ SLIDE 2 — Problem ==============================
s = slide(); title_bar(s, "The problem", "An AI answer is not verifiable evidence")
# big equation
box(s, 0.9, 2.15, 3.4, 1.5, fill=ACC_LT, ln=ACC, lnw=1.25)
txt(s, 0.9, 2.15, 3.4, 1.5, [[("AI ANSWER", 15, ACC, True)],[("“Approve the loan.”", 13, INK2, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.2)
txt(s, 4.5, 2.15, 1.2, 1.5, [[("≠", 54, RED, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 5.9, 2.15, 3.4, 1.5, fill=GREEN_LT, ln=GREEN, lnw=1.25)
txt(s, 5.9, 2.15, 3.4, 1.5, [[("VERIFIABLE EVIDENCE", 15, GREEN, True)],[("checkable by someone else", 13, INK2, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.2)
txt(s, 0.9, 3.95, 8.4, 0.5, [[("Ordinary logs are mutable and platform-specific. A model’s explanation is generated prose, not a record.", 13, INK2, False)]])
# four gaps
gaps = [("Source provenance", "what did it rely on?"),
        ("Tool / model actions", "what did it actually do?"),
        ("Tamper detection", "was it changed after the fact?"),
        ("Independent check", "can someone else verify it, offline?")]
for i,(a,b) in enumerate(gaps):
    x = 0.9 + i*3.0
    box(s, x, 4.75, 2.75, 1.7, fill=None, ln=LINE, lnw=1.25)
    rule(s, x+0.22, 5.02, 0.5, col=ACC, h=0.05)
    txt(s, x+0.22, 5.16, 2.35, 1.2, [[(a, 14, INK, True)],[(b, 12, INK2, False)]], line=1.15)
footer(s, 2)

# ============================ SLIDE 3 — Evolution ============================
s = slide(); title_bar(s, "Project evolution", "From more opinions to verifiable evidence")
stages = [("Orallexa", "multi-agent debate\n(bull · bear · judge)", ACC_LT, ACC),
          ("The lesson", "more opinions\n≠ more trust", AMBER_LT, AMBER),
          ("Shadow", "portable, verifiable\nevidence layer", GREEN_LT, GREEN)]
for i,(t,d,lt,c) in enumerate(stages):
    x = 0.9 + i*4.05
    box(s, x, 2.7, 3.5, 1.9, fill=lt, ln=c, lnw=1.5)
    txt(s, x, 2.9, 3.5, 0.5, [[(t, 20, c, True)]], align=PP_ALIGN.CENTER)
    txt(s, x, 3.5, 3.5, 1.0, [[(l, 14, INK2, False)] for l in d.split("\n")], align=PP_ALIGN.CENTER, line=1.2)
    if i < 2:
        txt(s, x+3.5, 2.7, 0.55, 1.9, [[("→", 34, MUTE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.9, 5.1, 11.5, 1.0,
    [[("Trust moved ", 16, INK, False),("down a layer", 16, ACC, True),
      (" — to the evidence underneath, not the voices on top.", 16, INK, False)],
     [("A refinement of the original research question, not a project that failed.", 13, INK2, False)]], line=1.4)
footer(s, 3)

# ============================ SLIDE 4 — Lifecycle ============================
s = slide(); title_bar(s, "What Shadow records", "The evidence lifecycle")
nodes = ["Source","Agent /\nTool Action","Evidence\nEvent","Hash\nChain","Signature","Independent\nVerification"]
cols  = [INK2, INK2, ACC, ACC, ACC, GREEN]
n = len(nodes); w = 1.78; gap = (11.9 - n*w)/(n-1)
for i,(lab,c) in enumerate(zip(nodes,cols)):
    x = 0.72 + i*(w+gap)
    box(s, x, 2.9, w, 1.35, fill=(GREEN_LT if i==n-1 else ACC_LT if i>=2 else RGBColor(0xF4,0xF4,0xF0)),
        ln=c, lnw=1.5)
    txt(s, x, 2.9, w, 1.35, [[(l, 13, c, True)] for l in lab.split("\n")], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.05)
    if i < n-1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, I(x+w+0.02), I(3.42), I(gap-0.04), I(0.32)); _noshadow(ar)
        ar.fill.solid(); ar.fill.fore_color.rgb = LINE; ar.line.fill.background()
txt(s, 0.9, 4.75, 11.5, 1.3,
    [[("•  Each event is structured and canonically serialized.", 15, INK, False)],
     [("•  Each event commits to the previous one — change an early event, everything after it ", 15, INK, False),("breaks", 15, RED, True),(".", 15, INK, False)],
     [("•  The sealed batch root is signed (Ed25519). Anyone can verify the record — ", 15, INK, False),("including offline", 15, ACC, True),(".", 15, INK, False)]], line=1.5)
footer(s, 4)

# ============================ SLIDE 5 — Architecture =========================
s = slide(); title_bar(s, "Architecture", "One core, three profiles, many surfaces")
def band(y, label, items, lt, c, note):
    box(s, 0.9, y, 11.5, 1.35, fill=lt, ln=c, lnw=1.25)
    txt(s, 1.15, y+0.12, 2.7, 1.1, [[(label, 16, c, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 3.9, y+0.14, 7.2, 1.1, [[(items, 13.5, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, line=1.25)
    txt(s, 3.9, y+0.98, 8.3, 0.3, [[(note, 10.5, MUTE, True)]])
band(2.1, "CORE", "canonical evidence schema  ·  signed hash-chain record  ·  source maps\nverifier  ·  claim–evidence graph  ·  ingested-output audit (structural)", ACC_LT, ACC, "host-tested")
band(3.65, "PROFILES", "banking-v1     ·     data-science-v1     ·     coding-agent-v1", RGBColor(0xF4,0xF4,0xF0), INK2, "one verification grammar across domains")
band(5.2, "INTERFACES", "CLI  ·  MCP (11 tools)  ·  HTTP  ·  browser verify.html   |   Unity Shadow Lens  ·  Three.js replay  ·  Android mock", GREEN_LT, GREEN, "verifier: browser-rendered  ·  Unity/3D: prototypes, device-validation pending")
txt(s, 0.9, 6.62, 11.5, 0.4, [[("Interfaces are ", 12, INK2, False),("not", 12, RED, True),(" equally mature — the deck labels which is which.", 12, INK2, False)]])
footer(s, 5)

# ============================ SLIDE 6 — Profiles =============================
s = slide(); title_bar(s, "Three profiles, one grammar", "Same machinery across domains")
rows = [("Banking", "document → risk claim → source", ACC),
        ("Data Science", "dataset → model → metric", ACC),
        ("Coding Agent", "issue → diff → tests → commit", ACC)]
for i,(a,b,c) in enumerate(rows):
    y = 2.2 + i*1.15
    box(s, 0.9, y, 3.2, 0.95, fill=ACC_LT, ln=ACC, lnw=1.25)
    txt(s, 0.9, y, 3.2, 0.95, [[(a, 17, ACC, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.3, y, 8.1, 0.95, fill=None, ln=LINE, lnw=1.25)
    txt(s, 4.55, y, 7.7, 0.95, [[(b, 16, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.9, 5.85, 11.5, 0.95, fill=RGBColor(0xF4,0xF4,0xF0), ln=None)
txt(s, 1.15, 5.85, 11.0, 0.95, [[("Underneath, the same verification grammar: ", 15, INK, False),
    ("sequence → hash chain → signature → source resolution.", 15, ACC, True)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# ============================ SLIDE 7 — TAMPER (hero) ========================
s = slide()
txt(s, 0.7, 0.5, 11.9, 0.4, [[("THE KEY MOMENT", 12, RED, True)]])
txt(s, 0.7, 0.82, 11.9, 0.9, [[("Tamper demonstration", 30, INK, True)]])
rule(s, 0.72, 1.62, 1.5, col=RED)
# pristine chain (green)
txt(s, 0.9, 1.95, 4, 0.35, [[("PRISTINE", 13, GREEN, True)]])
for i in range(6):
    x = 0.9 + i*1.95
    box(s, x, 2.35, 1.6, 0.75, fill=GREEN_LT, ln=GREEN, lnw=1.5)
    txt(s, x, 2.35, 1.6, 0.75, [[(f"seq {i}", 13, GREEN, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i<5: box(s, x+1.6, 2.68, 0.35, 0.08, fill=GREEN, round_=False)
# tampered chain
txt(s, 0.9, 3.45, 8, 0.35, [[("TAMPERED — one field changed at seq 2", 13, RED, True)]])
states = [("seq 0",GREEN),("seq 1",GREEN),("seq 2",RED),("seq 3",AMBER),("seq 4",AMBER),("seq 5",AMBER)]
for i,(lab,c) in enumerate(states):
    x = 0.9 + i*1.95
    lt = GREEN_LT if c==GREEN else RED_LT if c==RED else AMBER_LT
    box(s, x, 3.85, 1.6, 0.75, fill=lt, ln=c, lnw=1.5)
    tag = lab if c!=RED else lab+" ✗"
    txt(s, x, 3.85, 1.6, 0.75, [[(tag, 13, c, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i<5:
        col = GREEN if (i<1) else RED
        box(s, x+1.6, 4.18, 0.35, 0.08, fill=col, round_=False)
box(s, 5.72, 3.5, 3.0, 0.5, fill=None, ln=None)
txt(s, 4.6, 3.5, 2.1, 0.4, [[("↑ first failed sequence", 12, RED, True)]])
# takeaways
box(s, 0.9, 4.95, 11.5, 1.55, fill=RGBColor(0xF4,0xF4,0xF0), ln=None)
txt(s, 1.2, 5.1, 11.0, 1.4,
    [[("The verifier names the ", 15, INK, False),("exact first failed sequence", 15, RED, True),(" — not “something is wrong.”", 15, INK, False)],
     [("Everything after it is ", 15, INK, False),("invalidated", 15, RED, True),(" — the signature covered a chain that no longer exists.", 15, INK, False)],
     [("Even in failure, integrity and ", 15, INK, False),("correctness stay separate", 15, ACC, True),(" — it never says the conclusion was right.", 15, INK, False)]], line=1.45)
footer(s, 7)

# ============================ SLIDE 8 — Verify the Verifier ==================
s = slide(); title_bar(s, "Verify the Verifier", "Can you trust the verifier page itself?")
img = os.path.join(FIG, "fig07-verify-the-verifier-en.png")
if os.path.exists(img):
    box(s, 6.55, 2.0, 6.0, 3.9, fill=None, ln=LINE, lnw=1.0)
    try: s.shapes.add_picture(img, I(6.65), I(2.1), width=I(5.8))
    except Exception: pass
    txt(s, 6.55, 5.9, 6.0, 0.35, [[("Real render — Verify-the-Verifier (fixture manifest)", 10, MUTE, True)]])
pts = [("It carries a signed fixture manifest — ", "ASSETS MATCH SIGNED MANIFEST", GREEN),
       ("Real trust needs an out-of-band check: asset hashes vs. a manifest from a separate channel, and the release-key fingerprint compared independently — ", "INDEPENDENT COMPARISON NOT PERFORMED", AMBER),
       ("A page hashing itself is ", "not trust", RED)]
y=2.2
for a,b,c in pts:
    box(s, 0.9, y, 0.09, 1.05, fill=c, round_=False)
    txt(s, 1.15, y, 5.1, 1.15, [[(a, 13.5, INK, False),(b, 13.5, c, True),(".", 13.5, INK, False)]], line=1.25)
    y += 1.25
txt(s, 0.9, 6.15, 5.4, 0.5, [[("Current signing is ", 13, INK2, False),("fixture, not production.", 13, AMBER, True)]])
footer(s, 8)

# ============================ SLIDE 9 — Spatial =============================
s = slide(); title_bar(s, "The spatial experience", "3D for sequence, provenance, and tamper replay")
box(s, 0.9, 2.15, 5.5, 4.15, fill=ACC_LT, ln=ACC, lnw=1.25)
txt(s, 1.2, 2.4, 4.9, 3.7,
    [[("Unity Shadow Lens + Three.js replay", 16, ACC, True)],
     [("", 6, INK, False)],
     [("•  Three workspaces; provenance audit arc; Verify / Tamper / Reset", 14, INK, False)],
     [("•  Head-directed focus (gaze hover)", 14, INK, False)],
     [("•  A precise 2D fallback always remains", 14, INK, False)],
     [("•  Android mock APK built — 24.4 MB, hashed", 14, INK, False)]], line=1.3)
box(s, 6.7, 2.15, 5.7, 4.15, fill=AMBER_LT, ln=AMBER, lnw=1.25)
txt(s, 7.0, 2.4, 5.1, 3.7,
    [[("Stated honestly", 16, AMBER, True)],
     [("", 6, INK, False)],
     [("✗  NOT eye tracking — hover only", 14, INK, False)],
     [("✗  No RGB capture", 14, INK, False)],
     [("✗  No 6DoF on the mock build", 14, INK, False)],
     [("○  Android APK is built, ", 14, INK, False),("not device-validated", 14, AMBER, True)],
     [("○  Beam Pro / XREAL validation pending", 14, INK, False)],
     [("○  No user study yet — 3D benefit unproven", 14, INK, False)]], line=1.3)
footer(s, 9)

# ============================ SLIDE 10 — Evaluation =========================
s = slide(); title_bar(s, "Evaluation", "What is real today, measured honestly")
tiles = [("1,824 / 1,827", "host tests pass\n(3 skipped · 0 failed)", GREEN, GREEN_LT),
         ("Chromium 149", "browser verifier\nEN + 简体中文 · 0 external", ACC, ACC_LT),
         ("24.4 MB APK", "Android mock built\n(not device-validated)", AMBER, AMBER_LT),
         ("No user study", "3D comprehension\nnot yet proven (RQ4)", INK2, RGBColor(0xF1,0xF1,0xEC))]
for i,(big,sub,c,lt) in enumerate(tiles):
    x = 0.9 + i*2.95
    box(s, x, 2.4, 2.7, 2.3, fill=lt, ln=c, lnw=1.25)
    txt(s, x, 2.7, 2.7, 0.7, [[(big, 20, c, True)]], align=PP_ALIGN.CENTER)
    txt(s, x+0.15, 3.5, 2.4, 1.1, [[(l, 12.5, INK, False)] for l in sub.split("\n")], align=PP_ALIGN.CENTER, line=1.2)
txt(s, 0.9, 5.15, 11.5, 1.2,
    [[("Browser acceptance run in real Chromium (Playwright): 8 flows, English + Simplified Chinese,", 13.5, INK2, False)],
     [("CSP with 0 external requests / 0 violations, offline verification — real screenshots, not mock-ups.", 13.5, INK2, False)]], line=1.4)
footer(s, 10)

# ============================ SLIDE 11 — Where it stands ====================
s = slide(); title_bar(s, "Where it stands", "Implemented now vs. Capstone II")
box(s, 0.9, 2.2, 5.6, 4.1, fill=GREEN_LT, ln=GREEN, lnw=1.25)
txt(s, 1.15, 2.35, 5.1, 0.5, [[("IMPLEMENTED NOW", 15, GREEN, True)]])
for i,t in enumerate(["Deterministic evidence + hash chain","Independent + offline verification",
                      "Exact tamper localization","Bilingual verifier (EN / zh-CN)","Android mock build"]):
    txt(s, 1.15, 2.95+i*0.62, 5.1, 0.6, [[("✓  ", 14, GREEN, True),(t, 14, INK, False)]])
box(s, 6.85, 2.2, 5.55, 4.1, fill=AMBER_LT, ln=AMBER, lnw=1.25)
txt(s, 7.1, 2.35, 5.1, 0.5, [[("PENDING — CAPSTONE II", 15, AMBER, True)]])
for i,t in enumerate(["Beam Pro / XREAL device validation","Shared Unity / Three.js scene contract",
                      "User study — does 3D actually help?","Semantic audit of ingested output","Production signing · KMS / HSM"]):
    txt(s, 7.1, 2.95+i*0.62, 5.1, 0.6, [[("○  ", 14, AMBER, True),(t, 14, INK, False)]])
txt(s, 0.9, 6.55, 11.5, 0.4, [[("Fixture-signed · device-pending · no correctness claim — stated plainly, not oversold.", 12.5, INK2, True)]])
footer(s, 11)

# ============================ SLIDE 12 — Contribution =======================
s = slide()
box(s, 0, 0, SW, SH, fill=DARK, round_=False)
rule(s, 0, 0, SW, col=ACC, h=0.12)
txt(s, 0.9, 1.5, 11.5, 0.5, [[("CONTRIBUTION", 14, RGBColor(0x9F,0xB6,0xD6), True)]])
txt(s, 0.9, 2.25, 11.6, 2.6,
    [[("Shadow does not ask the auditor to", 30, RGBColor(0xE7,0xEC,0xF5), False)],
     [("trust the AI’s answer.", 30, RGBColor(0xE7,0xEC,0xF5), False)],
     [("", 12, INK, False)],
     [("It gives them evidence they can", 30, RGBColor(0xFF,0xFF,0xFF), True)],
     [("independently verify.", 30, RGBColor(0xFF,0xFF,0xFF), True)]], line=1.15)
rule(s, 0.92, 5.35, 2.0, col=ACC)
txt(s, 0.9, 5.6, 11.6, 0.9,
    [[("The sequence, the sources, the signature — and, if it was tampered with, the exact break.", 15, RGBColor(0x9F,0xB6,0xD6), False)],
     [("Integrity, not correctness.", 17, RGBColor(0x6F,0xD1,0x9A), True)]], line=1.35)
txt(s, 0.9, 6.8, 11.6, 0.4, [[("Thank you — questions welcome.", 14, RGBColor(0x8A,0x9C,0xBE), False)]])

out = os.path.join(HERE, "SHADOW_CAPSTONE_PRACTICE_PRESENTATION.pptx")
prs.save(out)
print(f"saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
